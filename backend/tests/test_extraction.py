import io

import docx
import pytest
from pptx import Presentation

from modules.extraction import (
    UnsafeURLError,
    UnsupportedFileTypeError,
    extract_from_file,
    extract_from_raw_text,
    extract_from_url,
)
from modules.extraction.docx_extractor import extract_docx_text
from modules.extraction.pdf_extractor import extract_pdf_text
from modules.extraction.pptx_extractor import extract_pptx_text
from modules.extraction.txt_extractor import extract_txt_text


def _build_docx_bytes(paragraphs):
    document = docx.Document()
    for text in paragraphs:
        document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _build_pptx_bytes(slide_texts):
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    for text in slide_texts:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_extract_docx_text_returns_paragraphs():
    data = _build_docx_bytes(["Distributed systems ensure consistency.", "Replication enables scaling."])
    text = extract_docx_text(data)
    assert "Distributed systems ensure consistency." in text
    assert "Replication enables scaling." in text


def test_extract_pptx_text_returns_slide_titles():
    data = _build_pptx_bytes(["Vector Databases", "Semantic Search"])
    text = extract_pptx_text(data)
    assert "Vector Databases" in text
    assert "Semantic Search" in text


def test_extract_txt_text_decodes_utf8():
    assert extract_txt_text("hello world".encode("utf-8")) == "hello world"


def test_extract_pdf_text_uses_pdfplumber(monkeypatch):
    class FakePage:
        def extract_text(self):
            return "Page one content."

    class FakePdf:
        pages = [FakePage()]

        def __enter__(self):
            return self

        def __exit__(self, *args):
            return False

    monkeypatch.setattr("modules.extraction.pdf_extractor.pdfplumber.open", lambda _f: FakePdf())
    assert extract_pdf_text(b"irrelevant") == "Page one content."


def test_extract_from_file_dispatches_by_extension():
    data = _build_docx_bytes(["Some content."])
    text = extract_from_file("notes.docx", data)
    assert "Some content." in text


def test_extract_from_file_rejects_unsupported_extension():
    with pytest.raises(UnsupportedFileTypeError):
        extract_from_file("archive.zip", b"data")


def test_extract_from_raw_text_passthrough():
    assert extract_from_raw_text("raw input") == "raw input"


def test_extract_from_url_rejects_private_addresses():
    with pytest.raises(UnsafeURLError):
        extract_from_url("http://127.0.0.1/secret")


def test_extract_from_url_rejects_non_http_scheme():
    with pytest.raises(UnsafeURLError):
        extract_from_url("file:///etc/passwd")
