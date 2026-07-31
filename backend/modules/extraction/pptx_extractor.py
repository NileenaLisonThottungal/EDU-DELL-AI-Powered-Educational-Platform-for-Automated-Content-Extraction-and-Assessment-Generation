"""PPTX text extraction (Section 3.2.1: python-pptx, titles/content/bullets)."""
import io

from pptx import Presentation


def extract_pptx_text(file_bytes: bytes) -> str:
    presentation = Presentation(io.BytesIO(file_bytes))
    chunks = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if text.strip():
                    chunks.append(text)

    return "\n".join(chunks)
